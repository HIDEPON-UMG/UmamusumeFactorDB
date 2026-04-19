"""UMG因子保管庫 システム構成図 (PowerPoint) を生成する。

使い方:
    python scripts/generate_architecture_pptx.py

出力:
    docs/system_architecture.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt


# カラーパレット
COLOR_USER = RGBColor(0x37, 0x4C, 0x80)        # 紺
COLOR_FORM = RGBColor(0x2D, 0x7D, 0xFF)        # 青
COLOR_GAS = RGBColor(0x0F, 0x9D, 0x58)         # 緑（Google）
COLOR_SHEET = RGBColor(0x34, 0xA8, 0x53)       # Sheets 緑
COLOR_DRIVE = RGBColor(0xFB, 0xBC, 0x05)       # Drive 黄
COLOR_CLOUDRUN = RGBColor(0x42, 0x85, 0xF4)    # Cloud Run 青
COLOR_SECRET = RGBColor(0xEA, 0x43, 0x35)      # Secret 赤
COLOR_DISCORD = RGBColor(0x58, 0x65, 0xF2)     # Discord 紫
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK = RGBColor(0x20, 0x20, 0x20)
COLOR_LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
COLOR_GRAY = RGBColor(0x7B, 0x84, 0x8F)


def _add_box(slide, x, y, w, h, fill, text, font_size=11, font_color=COLOR_WHITE, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.margin_left = Emu(80000)
    tf.margin_right = Emu(80000)
    tf.margin_top = Emu(50000)
    tf.margin_bottom = Emu(50000)
    tf.word_wrap = True
    tf.text = text
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = font_color
            run.font.name = "Meiryo UI"
    return shape


def _add_label(slide, x, y, w, h, text, font_size=10, font_color=COLOR_DARK, bold=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.LEFT
        for run in para.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = font_color
            run.font.name = "Meiryo UI"
    return tb


def _add_arrow(slide, x1, y1, x2, y2, color=COLOR_GRAY, label=None, label_offset=(0, 0)):
    line = slide.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = straight connector
    line.line.color.rgb = color
    line.line.width = Emu(20000)
    # 矢印の end タイプ
    ln = line.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle"})
    ln.append(tail)
    if label:
        tb = slide.shapes.add_textbox(
            (x1 + x2) // 2 + Emu(label_offset[0]),
            (y1 + y2) // 2 + Emu(label_offset[1]),
            Inches(2.2),
            Inches(0.3),
        )
        tf = tb.text_frame
        tf.text = label
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.LEFT
            for run in para.runs:
                run.font.size = Pt(9)
                run.font.color.rgb = COLOR_DARK
                run.font.name = "Meiryo UI"
    return line


def _add_title(slide, text, subtitle=None):
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12), Inches(0.8))
    tf = tb.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    for run in p.runs:
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = COLOR_DARK
        run.font.name = "Meiryo UI"
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        for run in p2.runs:
            run.font.size = Pt(12)
            run.font.color.rgb = COLOR_GRAY
            run.font.name = "Meiryo UI"


def _add_section_label(slide, x, y, w, h, text):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.text = text
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = COLOR_GRAY
            run.font.name = "Meiryo UI"


def build_overview_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_title(slide, "UMG因子保管庫 システム構成", "投稿フロー + 検索フロー + Discord 通知")

    # 左端：ユーザー
    _add_box(slide, Inches(0.3), Inches(3.0), Inches(1.4), Inches(0.8), COLOR_USER,
             "👤 ユーザー\n（トレーナー）", font_size=11)

    # Google Form
    _add_box(slide, Inches(2.0), Inches(1.0), Inches(2.0), Inches(0.9), COLOR_FORM,
             "Google Form\n投稿フォーム", font_size=12)

    # 検索ブラウザ
    _add_box(slide, Inches(2.0), Inches(5.3), Inches(2.0), Inches(0.9), COLOR_FORM,
             "ブラウザ\n検索 UI (search.html)", font_size=12)

    # Apps Script（中央）
    _add_box(slide, Inches(5.0), Inches(3.0), Inches(3.0), Inches(1.6), COLOR_GAS,
             "Apps Script (Code.gs)\n・onFormSubmit\n・doPost (webhook)\n・doGet (検索 API)\n・Discord 通知", font_size=11)

    # Google Sheets
    _add_box(slide, Inches(5.0), Inches(5.4), Inches(3.0), Inches(0.9), COLOR_SHEET,
             "Google Sheets\nfactors_normalized / 応答タブ", font_size=11)

    # Google Drive
    _add_box(slide, Inches(5.0), Inches(1.0), Inches(3.0), Inches(0.9), COLOR_DRIVE,
             "Google Drive\n因子画像（匿名リネーム済）", font_size=11)

    # Cloud Run
    _add_box(slide, Inches(9.0), Inches(2.2), Inches(3.8), Inches(1.5), COLOR_CLOUDRUN,
             "Cloud Run (factor-processor)\nFastAPI + ONNX + EasyOCR\n・因子解析\n・Apps Script へ結果 POST", font_size=11)

    # Secret Manager
    _add_box(slide, Inches(9.0), Inches(4.0), Inches(3.8), Inches(0.8), COLOR_SECRET,
             "GCP Secret Manager\napps-script-secret / cloud-run-shared-secret", font_size=10)

    # Discord
    _add_box(slide, Inches(9.0), Inches(5.3), Inches(3.8), Inches(1.3), COLOR_DISCORD,
             "Discord Webhook × 3\n対人 / 査定 / 競技場\n（キタサンブラック bot）", font_size=11)

    # 矢印：投稿フロー
    _add_arrow(slide, Inches(1.7), Inches(3.2), Inches(2.0), Inches(1.5),
               label="① 画像+連絡先+目的", label_offset=(-100000, -200000))
    _add_arrow(slide, Inches(4.0), Inches(1.4), Inches(5.0), Inches(1.4),
               label="② 保存", label_offset=(0, -250000))
    _add_arrow(slide, Inches(6.5), Inches(1.9), Inches(6.5), Inches(3.0),
               label="③ Driveから取得", label_offset=(100000, 0))
    _add_arrow(slide, Inches(8.0), Inches(3.5), Inches(9.0), Inches(2.8),
               label="④ base64+secret", label_offset=(0, -200000))
    _add_arrow(slide, Inches(9.0), Inches(3.5), Inches(8.0), Inches(3.8),
               label="⑤ 解析結果", label_offset=(-1400000, 100000))
    _add_arrow(slide, Inches(6.5), Inches(4.6), Inches(6.5), Inches(5.4),
               label="⑥ 3行追記", label_offset=(100000, -100000))
    _add_arrow(slide, Inches(8.0), Inches(4.0), Inches(9.0), Inches(5.6),
               label="⑦ Discord通知", label_offset=(-100000, -200000))

    # 矢印：検索フロー（破線風）
    _add_arrow(slide, Inches(1.7), Inches(5.7), Inches(2.0), Inches(5.7),
               color=RGBColor(0xFF, 0x8A, 0x65), label="A 開く", label_offset=(-300000, -250000))
    _add_arrow(slide, Inches(4.0), Inches(5.7), Inches(5.0), Inches(4.0),
               color=RGBColor(0xFF, 0x8A, 0x65), label="B google.script.run", label_offset=(0, 100000))
    _add_arrow(slide, Inches(6.5), Inches(4.6), Inches(6.5), Inches(5.4),
               color=RGBColor(0xFF, 0x8A, 0x65))

    # 凡例
    _add_label(slide, Inches(0.4), Inches(6.7), Inches(12), Inches(0.4),
               "▶ グレー矢印：投稿フロー  　▶ オレンジ矢印：検索フロー",
               font_size=10, font_color=COLOR_GRAY)


def build_dataflow_slide(prs: Presentation):
    """投稿フロー詳細（シーケンス的）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "投稿フロー詳細", "Form 送信 → 解析 → スプレ → Discord 通知")

    # 縦軸に沿ってステップを箱で表現
    steps = [
        ("① ユーザー", "Google Form に画像 + 連絡先 + 目的・用途を送信", COLOR_USER),
        ("② Form → 応答タブ", "Google Sheets『フォームの回答 1』に自動記録", COLOR_FORM),
        ("③ onFormSubmit 発火", "Apps Script が応答行を読み、Drive ファイル ID 抽出", COLOR_GAS),
        ("④ 画像匿名リネーム", "factor_yyyyMMdd_HHmmss.ext 形式で setName", COLOR_DRIVE),
        ("⑤ Cloud Run 呼び出し", "POST /process に画像 base64 + secret を送信", COLOR_CLOUDRUN),
        ("⑥ ONNX + OCR 解析", "因子ボックス検出 → 青/赤/緑/白 + ★数 + character 逆引き", COLOR_CLOUDRUN),
        ("⑦ Cloud Run → webhook", "Apps Script doPost に 3 行の解析結果を POST", COLOR_CLOUDRUN),
        ("⑧ factors_normalized 書き込み", "submission_id 共通で main/parent1/parent2 の 3 行追加", COLOR_SHEET),
        ("⑨ 応答タブに status 書き戻し", "submission_id + processed を応答行へ", COLOR_SHEET),
        ("⑩ Discord 通知（対人/査定/競技場）", "画像公開化 → キタサン口調メッセージ + embed", COLOR_DISCORD),
    ]

    x = Inches(0.5)
    y_start = Inches(1.2)
    row_h = Inches(0.55)
    gap = Inches(0.05)
    for i, (title, detail, color) in enumerate(steps):
        y = y_start + (row_h + gap) * i
        _add_box(slide, x, y, Inches(4.5), row_h, color, title, font_size=11)
        _add_label(slide, x + Inches(4.7), y + Emu(40000), Inches(8.3), row_h,
                   detail, font_size=10, font_color=COLOR_DARK)


def build_data_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "データスキーマ", "factors_normalized タブ（1 投稿 = 3 行）")

    headers = ["列名", "型", "説明"]
    rows = [
        ("submission_id", "string", "投稿の一意 ID（3 行で共通）"),
        ("submitted_at", "ISO8601", "投稿日時"),
        ("submitter_id", "string", "連絡先（Discord 名など、任意）"),
        ("image_filename", "string", "サーバ内仮名"),
        ("role", "string", "main / parent1 / parent2"),
        ("character", "string", "[衣装名]キャラ名（固有スキル逆引きで確定）"),
        ("blue_type / blue_star", "string/int", "青因子 + ★数"),
        ("red_type / red_star", "string/int", "赤因子 + ★数"),
        ("green_name / green_star", "string/int", "緑因子（固有スキル） + ★数"),
        ("factor_01..60_name/star", "string/int", "白因子スロット（最大 60 件）"),
    ]
    n = len(rows) + 1
    top = Inches(1.3)
    left = Inches(0.5)
    col_w = [Inches(3.0), Inches(1.8), Inches(8.0)]
    row_h_header = Inches(0.35)
    row_h = Inches(0.32)

    # ヘッダー
    x = left
    for i, h in enumerate(headers):
        _add_box(slide, x, top, col_w[i], row_h_header, COLOR_DARK, h, font_size=10, font_color=COLOR_WHITE)
        x += col_w[i]

    for r_i, row in enumerate(rows):
        y = top + row_h_header + (row_h * r_i)
        x = left
        for c_i, cell in enumerate(row):
            color = COLOR_LIGHT if r_i % 2 == 0 else COLOR_WHITE
            _add_box(slide, x, y, col_w[c_i], row_h, color, cell,
                     font_size=9, font_color=COLOR_DARK, bold=False)
            x += col_w[c_i]

    _add_label(slide, Inches(0.5), Inches(5.6), Inches(12), Inches(0.4),
               "※ 1 投稿につき main / parent1 / parent2 の 3 行が submission_id 共通で並ぶ",
               font_size=10, font_color=COLOR_GRAY)
    _add_label(slide, Inches(0.5), Inches(5.95), Inches(12), Inches(0.4),
               "※ 目的・用途・画像 URL は『フォームの回答 1』タブ側に保存（search.html で submission_id 突合）",
               font_size=10, font_color=COLOR_GRAY)


def build_security_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "セキュリティ & 運用上の配慮")

    items = [
        ("🔒 画像ファイル名の匿名化", "Apps Script で factor_yyyyMMdd_HHmmss にリネーム（Google アカウント名の除去）", COLOR_SECRET),
        ("🔒 Webhook 認証", "Apps Script ↔ Cloud Run 間は SHARED_SECRET（Secret Manager で管理）", COLOR_SECRET),
        ("🔒 clickjacking 対策", "search.html は XFrameOptionsMode デフォルト（ALLOWALL 解除）", COLOR_SECRET),
        ("🔒 画像公開範囲", "Discord Webhook 送信対象のみ ANYONE_WITH_LINK（閲覧のみ）", COLOR_SECRET),
        ("🔒 XSS 対策", "search.html の全ユーザー入力は escapeHtml() を経由して描画", COLOR_SECRET),
        ("🛠 モデル更新", "scripts/fetch_unique_skills.py で UmaTools から固有スキル対応表を再生成", COLOR_GAS),
        ("🛠 Cloud Run 再デプロイ", "gcloud run deploy factor-processor --source . （新コード反映時）", COLOR_CLOUDRUN),
        ("🛠 Apps Script 再デプロイ", "コード保存後に『デプロイを管理 → 新バージョン』で反映", COLOR_GAS),
        ("⚠ タイムアウト", "Cloud Run 300 秒制限。コールドスタート 30〜60s + 解析 30〜60s を想定", COLOR_DRIVE),
        ("⚠ 個人情報", "連絡先（Discord 名）は検索画面で全閲覧者に見える仕様。投稿時の任意入力を推奨", COLOR_DRIVE),
    ]
    x = Inches(0.5)
    y_start = Inches(1.2)
    row_h = Inches(0.52)
    gap = Inches(0.05)
    for i, (title, detail, color) in enumerate(items):
        y = y_start + (row_h + gap) * i
        _add_box(slide, x, y, Inches(4.0), row_h, color, title, font_size=10)
        _add_label(slide, x + Inches(4.2), y + Emu(40000), Inches(8.8), row_h,
                   detail, font_size=9.5, font_color=COLOR_DARK)


def build_cover_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 背景帯
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.8), Inches(13.33), Inches(1.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_USER
    bar.line.color.rgb = COLOR_USER

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(12), Inches(0.8))
    tf = tb.text_frame
    tf.text = "UMG因子保管庫"
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = COLOR_WHITE
        run.font.name = "Meiryo UI"

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12), Inches(0.8))
    tf2 = tb2.text_frame
    tf2.text = "Umamusume Factor DB — システム構成資料"
    for run in tf2.paragraphs[0].runs:
        run.font.size = Pt(18)
        run.font.color.rgb = COLOR_WHITE
        run.font.name = "Meiryo UI"

    tb3 = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12), Inches(0.6))
    tf3 = tb3.text_frame
    tf3.text = "Google Form + Apps Script + Cloud Run + Google Sheets + Discord Webhook"
    for run in tf3.paragraphs[0].runs:
        run.font.size = Pt(12)
        run.font.color.rgb = COLOR_GRAY
        run.font.name = "Meiryo UI"


def main() -> int:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_cover_slide(prs)
    build_overview_slide(prs)
    build_dataflow_slide(prs)
    build_data_slide(prs)
    build_security_slide(prs)

    out_path = Path(__file__).resolve().parents[1] / "docs" / "system_architecture.pptx"
    prs.save(str(out_path))
    print(f"wrote {out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
